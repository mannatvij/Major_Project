"""
generate_final_ppt.py — Final end-term presentation.

Produces  Major_Project_Final_Presentation.pptx  in the repo root.
The earlier  generate_ppt.py / Major_Project_Presentation.pptx  is left untouched.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# ── Palette (matches OD deck) ───────────────────────────────────────────────
DARK_BLUE   = RGBColor(0x0D, 0x2C, 0x54)
MID_BLUE    = RGBColor(0x19, 0x76, 0xD2)
LIGHT_BLUE  = RGBColor(0xE3, 0xF2, 0xFD)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT      = RGBColor(0x00, 0xB0, 0x8A)
DARK_TEXT   = RGBColor(0x1A, 0x1A, 0x2E)
GREY_TEXT   = RGBColor(0x55, 0x55, 0x66)
TICK_GREEN  = RGBColor(0x2E, 0x7D, 0x32)
SUBTITLE    = RGBColor(0xBB, 0xDE, 0xFB)
PALE_BORDER = RGBColor(0xC5, 0xD8, 0xF0)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


# ── Helpers ─────────────────────────────────────────────────────────────────
def add_rect(slide, l, t, w, h, fill_rgb=None, line_rgb=None, line_width_pt=0):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.line.fill.background()
    if fill_rgb:
        s.fill.solid(); s.fill.fore_color.rgb = fill_rgb
    else:
        s.fill.background()
    if line_rgb:
        s.line.color.rgb = line_rgb
        s.line.width = Pt(line_width_pt)
    else:
        s.line.fill.background()
    return s


def add_triangle(slide, l, t, w, h, fill_rgb):
    """Right-pointing arrow connector."""
    s = slide.shapes.add_shape(5, Inches(l), Inches(t), Inches(w), Inches(h))  # 5 = right triangle
    s.line.fill.background()
    s.fill.solid(); s.fill.fore_color.rgb = fill_rgb
    return s


def add_textbox(slide, text, l, t, w, h,
                font_size=18, bold=False, color=WHITE,
                align=PP_ALIGN.LEFT, italic=False):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def header_bar(slide, title, subtitle=None):
    add_rect(slide, 0, 0, 13.33, 1.3, fill_rgb=DARK_BLUE)
    add_textbox(slide, title, 0.35, 0.12, 12.0, 0.7,
                font_size=32, bold=True, color=WHITE)
    if subtitle:
        add_textbox(slide, subtitle, 0.35, 0.72, 12.0, 0.45,
                    font_size=16, color=SUBTITLE)


def accent_line(slide):
    add_rect(slide, 0, 1.3, 13.33, 0.05, fill_rgb=ACCENT)


def slide_background(slide):
    add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=LIGHT_BLUE)


def footer(slide, text="Smart Healthcare Appointment System  |  Mannat Vij  |  Final End-Term Presentation"):
    add_rect(slide, 0, 7.18, 13.33, 0.32, fill_rgb=DARK_BLUE)
    add_textbox(slide, text, 0.3, 7.19, 12.7, 0.28,
                font_size=11, color=SUBTITLE, align=PP_ALIGN.CENTER)


def add_image(slide, path, l, t, w, h):
    if os.path.exists(path):
        slide.shapes.add_picture(path, Inches(l), Inches(t), Inches(w), Inches(h))
    else:
        add_rect(slide, l, t, w, h, fill_rgb=RGBColor(0xCC, 0xCC, 0xCC))
        add_textbox(slide, "[Screenshot pending]", l + 0.1, t + h / 2 - 0.2, w - 0.2, 0.4,
                    font_size=14, color=DARK_TEXT, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=DARK_BLUE)
add_rect(slide, 0, 0, 0.12, 7.5, fill_rgb=ACCENT)
add_rect(slide, 0, 6.8, 13.33, 0.7, fill_rgb=RGBColor(0x0A, 0x1F, 0x3D))
add_rect(slide, 0.5, 1.55, 12.3, 3.7, fill_rgb=RGBColor(0x12, 0x3A, 0x6B))

add_textbox(slide, "Smart Healthcare", 0.8, 1.7, 11.7, 0.9,
            font_size=54, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(slide, "Appointment System", 0.8, 2.55, 11.7, 0.9,
            font_size=54, bold=True, color=RGBColor(0x64, 0xB5, 0xF6), align=PP_ALIGN.CENTER)
add_rect(slide, 3.5, 3.55, 6.3, 0.06, fill_rgb=ACCENT)
add_textbox(slide, "FINAL PROJECT PRESENTATION  ·  END TERM", 0.8, 3.72, 11.7, 0.5,
            font_size=18, bold=True, color=SUBTITLE, align=PP_ALIGN.CENTER)
add_textbox(slide, "Presented by:  Mannat Vij", 0.8, 4.35, 11.7, 0.45,
            font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(slide, "Under the Guidance of:  Dr. Krishna Kumar", 0.8, 4.85, 11.7, 0.45,
            font_size=18, color=SUBTITLE, align=PP_ALIGN.CENTER)
add_textbox(slide, "Computer Science & Engineering  |  2024–25",
            0.8, 6.88, 11.7, 0.35,
            font_size=13, color=RGBColor(0x78, 0x9A, 0xBF), align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Index
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
slide_background(slide); header_bar(slide, "Index", "Presentation Overview"); accent_line(slide); footer(slide)

index_items = [
    ("01", "Introduction",          "Why this project & what it does"),
    ("02", "Objectives",            "Five goals of the system"),
    ("03", "Methodology",           "Approach, flow, tools used"),
    ("04", "Present Status & Features", "What is live and working today"),
    ("05", "Project Snapshots",     "Live screenshots of every module"),
    ("06", "Future Scope",          "What comes next"),
]
cols = [(0.45, 6.3), (6.9, 6.3)]
for idx, (num, title, desc) in enumerate(index_items):
    row = idx % 3; col = idx // 3
    lx, lw = cols[col]; ty = 1.55 + row * 1.7
    add_rect(slide, lx, ty, lw, 1.45, fill_rgb=WHITE, line_rgb=MID_BLUE, line_width_pt=1.2)
    add_rect(slide, lx + 0.12, ty + 0.18, 0.55, 0.55, fill_rgb=MID_BLUE)
    add_textbox(slide, num, lx + 0.12, ty + 0.18, 0.55, 0.55,
                font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, title, lx + 0.82, ty + 0.12, lw - 1.0, 0.42,
                font_size=20, bold=True, color=DARK_BLUE)
    add_textbox(slide, desc, lx + 0.82, ty + 0.55, lw - 1.0, 0.6,
                font_size=14, color=GREY_TEXT)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Introduction
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
slide_background(slide); header_bar(slide, "Introduction", "About the Project"); accent_line(slide); footer(slide)

intro_para = (
    "Healthcare today is held back by manual phone-bookings, paper prescriptions, "
    "missed reminders, and disconnected payment systems — patients struggle to find "
    "the right doctor and doctors lose time managing schedules and queues.\n\n"
    "The Smart Healthcare Appointment System unifies the entire patient journey on a "
    "single platform: an AI symptom-checker triages the patient, search and filter "
    "surface the right specialist, a visual slot-picker books the appointment, "
    "Razorpay handles online payment, Mailtrap delivers booking / cancellation / "
    "24-hour reminder emails, calendar files (.ics) sync to Google or Outlook, "
    "doctors write digital prescriptions that auto-email to the patient, and patients "
    "leave star-rated reviews. Admins oversee the platform with live KPI charts, "
    "doctor-approval workflows, system-health checks, and CSV exports.\n\n"
    "It is a full-stack web app — React 19 + Material-UI on the front, Spring Boot 3 "
    "+ MongoDB on the back, a Python Flask + scikit-learn microservice for ML triage "
    "(with a Java rule-based fallback), and a Playwright end-to-end suite that "
    "exercises the entire workflow."
)
add_rect(slide, 0.5, 1.55, 12.33, 5.4, fill_rgb=WHITE, line_rgb=PALE_BORDER, line_width_pt=1)
add_rect(slide, 0.5, 1.55, 0.18, 5.4, fill_rgb=MID_BLUE)
txb = slide.shapes.add_textbox(Inches(0.9), Inches(1.7), Inches(11.8), Inches(4.7))
txb.word_wrap = True
p = txb.text_frame.paragraphs[0]
p.alignment = PP_ALIGN.LEFT
r = p.add_run(); r.text = intro_para
r.font.size = Pt(15); r.font.color.rgb = DARK_TEXT

# Tech-stack chips
chips = ["React 19", "Spring Boot 3.2", "MongoDB", "JWT", "Material-UI",
         "Mailtrap SMTP", "Razorpay", "Flask + scikit-learn", "Recharts", "Playwright"]
chip_x = 0.9
for chip in chips:
    cw = len(chip) * 0.105 + 0.35
    add_rect(slide, chip_x, 6.55, cw, 0.38, fill_rgb=MID_BLUE)
    add_textbox(slide, chip, chip_x + 0.05, 6.56, cw - 0.1, 0.35,
                font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    chip_x += cw + 0.12


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Objectives
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
slide_background(slide); header_bar(slide, "Objectives", "Goals of the System"); accent_line(slide); footer(slide)

objectives = [
    ("01", "Unify the Full Healthcare Lifecycle on One Platform",
     "Register, triage, find a doctor, book a slot, pay online, attend, receive a digital prescription and "
     "leave a review — all in one web app, no external tools or paperwork."),
    ("02", "Three-Role Access with Scalable RESTful Backend",
     "Patient / Doctor / Admin portals over a Spring Boot 3 + MongoDB backend, documented live via Swagger UI "
     "as the single source of truth for every endpoint."),
    ("03", "AI-Assisted Symptom Triage",
     "Flask ML microservice (TF-IDF + RandomForest) recommends a specialisation through a multi-turn chat with "
     "red-flag detection; a Java rule-based fallback keeps chat alive when ML is offline."),
    ("04", "Reliable Communication & Calendaring",
     "Async Mailtrap email for booking / cancellation / 24-hour reminders + downloadable .ics file and Google / "
     "Outlook deep-links pre-filled with date and doctor info."),
    ("05", "Admin Oversight & Analytics",
     "KPI dashboard with Recharts (7/30/60-day trends), doctor-approval workflow, system-health panel, and "
     "one-click CSV exports for users, appointments, payments and reviews."),
]

for i, (num, title, desc) in enumerate(objectives):
    ty = 1.55 + i * 1.08
    add_rect(slide, 0.45, ty, 12.4, 0.98, fill_rgb=WHITE, line_rgb=PALE_BORDER, line_width_pt=1)
    add_rect(slide, 0.45, ty, 0.72, 0.98, fill_rgb=MID_BLUE)
    add_textbox(slide, num, 0.45, ty + 0.22, 0.72, 0.5,
                font_size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, title, 1.3, ty + 0.05, 11.2, 0.38,
                font_size=16.5, bold=True, color=DARK_BLUE)
    add_textbox(slide, desc, 1.3, ty + 0.43, 11.2, 0.55,
                font_size=12.5, color=GREY_TEXT)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Methodology (with flow chart)
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
slide_background(slide); header_bar(slide, "Methodology", "Development Approach, Flow & Tools"); accent_line(slide); footer(slide)

# Flow-chart strip
flow_steps = ["Plan", "Design DB & APIs", "Backend (Spring)", "Frontend (React)",
              "ML Service (Flask)", "Email & Payments", "Admin & Analytics", "E2E Test"]
flow_top = 1.5
flow_box_w = 1.45
flow_box_h = 0.6
arrow_w = 0.13
gap = 0.04
total = len(flow_steps) * flow_box_w + (len(flow_steps) - 1) * (arrow_w + 2 * gap)
start_x = (13.33 - total) / 2

for i, step in enumerate(flow_steps):
    bx = start_x + i * (flow_box_w + arrow_w + 2 * gap)
    add_rect(slide, bx, flow_top, flow_box_w, flow_box_h, fill_rgb=DARK_BLUE)
    add_textbox(slide, step, bx, flow_top, flow_box_w, flow_box_h,
                font_size=10.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if i < len(flow_steps) - 1:
        add_triangle(slide, bx + flow_box_w + gap, flow_top + 0.12, arrow_w, flow_box_h - 0.24,
                     fill_rgb=ACCENT)

# Step grid (8 cards in 2 columns, smaller to fit)
steps = [
    ("1", "Requirement Analysis",
     "Three roles (Patient / Doctor / Admin), use-cases, ER design, API contracts."),
    ("2", "Environment Setup",
     "Java 17 · Spring Boot 3.2 · React 19 · MongoDB · Node 18 · Python 3 / Flask."),
    ("3", "Database Design",
     "Collections: users, appointments, prescriptions, reviews, payments, reminders, chat_sessions."),
    ("4", "Backend REST APIs + Swagger UI",
     "Auth, Users, Doctors, Appointments, Prescriptions, Reviews, Payments — all live at /swagger-ui."),
    ("5", "Frontend SPA",
     "React with protected routes, Axios JWT interceptors, dark/light theme context, responsive layout."),
    ("6", "Email · Payments · Reminders",
     "Mailtrap SMTP with @Async; Razorpay checkout; @Scheduled cron at 09:00 daily for 24-hr reminders."),
    ("7", "ML Chatbot (with fallback)",
     "Flask multi-turn FSM, red-flag detection; Java SymptomMatcher activates if Flask is offline."),
    ("8", "Testing",
     "Postman + Swagger for API; Playwright single full-session E2E spec exercising every flow."),
]

step_top = 2.3
for col_idx, col_steps in enumerate([steps[:4], steps[4:]]):
    lx = 0.45 if col_idx == 0 else 6.88
    for row_i, (num, title, desc) in enumerate(col_steps):
        ty = step_top + row_i * 1.05
        add_rect(slide, lx, ty, 6.0, 0.95, fill_rgb=WHITE, line_rgb=PALE_BORDER, line_width_pt=1)
        add_rect(slide, lx, ty, 0.45, 0.95, fill_rgb=ACCENT)
        add_textbox(slide, num, lx, ty + 0.27, 0.45, 0.4,
                    font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, title, lx + 0.55, ty + 0.05, 5.4, 0.4,
                    font_size=13, bold=True, color=DARK_BLUE)
        add_textbox(slide, desc, lx + 0.55, ty + 0.4, 5.4, 0.55,
                    font_size=10.5, color=GREY_TEXT)

# Security badge
add_rect(slide, 3.0, 6.6, 7.33, 0.5, fill_rgb=DARK_BLUE)
add_textbox(slide, "Secured with JWT (HMAC-SHA256) · BCrypt password hashing · Spring Security role guards",
            3.0, 6.6, 7.33, 0.5, font_size=12, bold=True, color=SUBTITLE, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6A — Features (Patient)
# ════════════════════════════════════════════════════════════════════════════
def features_slide(title, subtitle, items):
    s = prs.slides.add_slide(BLANK)
    slide_background(s); header_bar(s, title, subtitle); accent_line(s); footer(s)
    cols = [items[: (len(items) + 1) // 2], items[(len(items) + 1) // 2:]]
    col_xs = [0.45, 6.88]
    for ci, col_items in enumerate(cols):
        lx = col_xs[ci]
        for ri, (name, desc) in enumerate(col_items):
            ty = 1.5 + ri * 0.92
            add_rect(s, lx, ty, 6.0, 0.84, fill_rgb=WHITE, line_rgb=PALE_BORDER, line_width_pt=1)
            add_rect(s, lx, ty, 0.55, 0.84, fill_rgb=TICK_GREEN)
            add_textbox(s, "✔", lx, ty + 0.16, 0.55, 0.5,
                        font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
            add_textbox(s, name, lx + 0.65, ty + 0.04, 5.3, 0.36,
                        font_size=12.5, bold=True, color=DARK_BLUE)
            add_textbox(s, desc, lx + 0.65, ty + 0.36, 5.3, 0.5,
                        font_size=10, color=GREY_TEXT)
    return s


patient_features = [
    ("Three-role JWT auth",
     "Separate signup/login flows for Patient, Doctor, Admin; HMAC-SHA256 tokens, BCrypt passwords."),
    ("AI Symptom Checker",
     "Multi-turn chat (symptom → severity → duration → red-flag) recommends a specialisation."),
    ("Doctor discovery",
     "Full-text name search + specialisation filter + paginated results + per-doctor profile pages."),
    ("Visual slot picker",
     "Live availability per doctor; only open slots are clickable in the booking form."),
    ("Razorpay online payment",
     "Pay before confirmation; payment record persisted and visible in patient history."),
    ("Calendar export",
     "One-click .ics download + Google / Outlook deep-links pre-filled with date & doctor."),
    ("Prescription viewer",
     "See every prescription from doctor: medications, dosage, frequency, advice, date."),
    ("Reviews & 5-star ratings",
     "Submit a rating + comment after a completed visit; aggregated on doctor profile."),
    ("Profile & password",
     "Edit personal info; change password with old-password verification."),
    ("Notification bell",
     "Live badge for booking confirmations, cancellations, reminders, prescriptions."),
    ("Dark / Light theme",
     "Toggle in app bar; preference saved to localStorage; respects OS preference initially."),
    ("Responsive UI",
     "Desktop sidebar + mobile bottom nav; Material-UI breakpoints throughout."),
]
features_slide("Present Status & Features — Part 1", "Patient-Side", patient_features)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6B — Features (Doctor / Admin / System)
# ════════════════════════════════════════════════════════════════════════════
admin_features = [
    ("Doctor schedule management",
     "Add or remove time slots dynamically; bulk weekly templates; one-click 'Out of Clinic'."),
    ("Doctor appointments queue",
     "Filter by status (Pending / Confirmed / Completed / Cancelled); accept / reject in a click."),
    ("Write prescription",
     "Structured form: medications, dosage, frequency, duration, advice — auto-emailed to patient."),
    ("Doctor profile page",
     "Public bio + cumulative rating + full review list visible to patients."),
    ("Admin dashboard",
     "KPI cards (users, doctors, appointments, revenue) + Recharts line / bar / pie charts."),
    ("Admin user management",
     "Searchable user table; activate, deactivate, delete; role and status filters."),
    ("Doctor approval workflow",
     "New doctors stay invisible until admin approves; rejection toggles active flag."),
    ("CSV exports",
     "One-click export for Users, Appointments, Payments, Reviews."),
    ("System-health panel",
     "Live indicators for Mongo, Mailer, ML service, Razorpay; live activity feed below."),
    ("Mailtrap email pipeline",
     "Async @Async-driven booking confirmation, cancellation notice, 24-hr reminder mails."),
    ("24-hour automated reminders",
     "@Scheduled cron at 09:00 daily; Reminders collection prevents duplicates."),
    ("Cancellation flow",
     "Slot auto-released back into availability + email sent to the patient."),
    ("Swagger UI live API docs",
     "/swagger-ui.html — single source of truth for every endpoint, interactively testable."),
    ("Playwright E2E suite",
     "Full-session spec: register → triage → book → pay → accept → prescribe → review → admin."),
]
features_slide("Present Status & Features — Part 2", "Doctor · Admin · System", admin_features)


# ════════════════════════════════════════════════════════════════════════════
# SLIDES 7–14 — Project Snapshots
# ════════════════════════════════════════════════════════════════════════════
snapshot_slides = [
    ("Authentication & Onboarding",
     ("final_login.png",          "Login Page"),
     ("final_register.png",       "Register Page")),
    ("AI Triage & Doctor Discovery",
     ("final_chatbot.png",        "AI Symptom Checker Chatbot"),
     ("final_doctor_list.png",    "Browse Doctors")),
    ("Doctor Profile & Booking",
     ("final_doctor_profile.png", "Doctor Profile + Reviews"),
     ("final_book_appointment.png","Booking · Slot Picker")),
    ("Payment & My Appointments",
     ("final_payment.png",        "Razorpay Online Payment"),
     ("final_appointments.png",   "Patient Appointments")),
    ("Notifications, Prescriptions & Dark Mode",
     ("final_prescription_view.png","Prescription Viewer"),
     ("final_dark_mode.png",      "Dark Mode")),
    ("Doctor Workspace",
     ("final_doctor_dashboard.png","Doctor Dashboard"),
     ("final_doctor_appointments.png","Doctor Appointments Queue")),
    ("Doctor Tools — Prescription & Availability",
     ("final_prescription_write.png","Writing a Prescription"),
     ("final_doctor_availability.png","Schedule Management")),
    ("Admin Portal & API Docs",
     ("final_admin_dashboard.png","Admin Dashboard · KPIs & Charts"),
     ("final_doctor_approvals.png","Doctor Approvals")),
    ("Admin Operations & Swagger",
     ("final_admin_user_mgmt.png","User Management"),
     ("final_swagger.png",        "Swagger UI · Live API Docs")),
]

for s_title, (lp, lcap), (rp, rcap) in snapshot_slides:
    slide = prs.slides.add_slide(BLANK)
    slide_background(slide); header_bar(slide, "Project Snapshots", s_title); accent_line(slide); footer(slide)
    for path, l, t, w, h, caption in [
        (lp, 0.45, 1.45, 5.9, 5.5, lcap),
        (rp, 6.95, 1.45, 5.9, 5.5, rcap),
    ]:
        add_rect(slide, l - 0.05, t - 0.05, w + 0.1, h + 0.1,
                 fill_rgb=WHITE, line_rgb=MID_BLUE, line_width_pt=1.5)
        add_image(slide, path, l, t, w, h)
        add_rect(slide, l, t + h, w, 0.35, fill_rgb=DARK_BLUE)
        add_textbox(slide, caption, l, t + h, w, 0.35,
                    font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE — Future Scope
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
slide_background(slide); header_bar(slide, "Future Scope", "Planned Enhancements"); accent_line(slide); footer(slide)

future_items = [
    ("Video Consultation (WebRTC)",
     "In-platform virtual visits with screen-share and optional session recording — patients consult without "
     "ever leaving the portal."),
    ("Mobile App (React Native)",
     "Native iOS / Android client with push notifications, offline-cached prescriptions, biometric login and "
     "native calendar sync."),
    ("Predictive Analytics & EHR Interop",
     "Forecast appointment volumes, predict no-shows, and export records to hospital EHRs via the FHIR "
     "standard for hospital-wide integration."),
]
for i, (title, desc) in enumerate(future_items):
    ty = 1.7 + i * 1.7
    add_rect(slide, 1.0, ty, 11.33, 1.5, fill_rgb=WHITE, line_rgb=PALE_BORDER, line_width_pt=1)
    add_rect(slide, 1.0, ty, 0.18, 1.5, fill_rgb=ACCENT)
    add_textbox(slide, title, 1.35, ty + 0.18, 10.8, 0.5,
                font_size=20, bold=True, color=DARK_BLUE)
    add_textbox(slide, desc, 1.35, ty + 0.7, 10.8, 0.78,
                font_size=14, color=GREY_TEXT)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE — Thank You
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=DARK_BLUE)
add_rect(slide, 0, 0, 0.12, 7.5, fill_rgb=ACCENT)
add_rect(slide, 0, 6.8, 13.33, 0.7, fill_rgb=RGBColor(0x0A, 0x1F, 0x3D))
add_rect(slide, 0.5, 2.2, 12.3, 2.8, fill_rgb=RGBColor(0x12, 0x3A, 0x6B))

add_textbox(slide, "Thank You", 0.8, 2.4, 11.7, 1.1,
            font_size=60, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_rect(slide, 4.0, 3.55, 5.3, 0.07, fill_rgb=ACCENT)
add_textbox(slide, "Smart Healthcare Appointment System", 0.8, 3.75, 11.7, 0.5,
            font_size=20, color=SUBTITLE, align=PP_ALIGN.CENTER)
add_textbox(slide, "Mannat Vij  |  Under the Guidance of Dr. Krishna Kumar",
            0.8, 4.3, 11.7, 0.4,
            font_size=16, color=RGBColor(0x9E, 0xC8, 0xF0), align=PP_ALIGN.CENTER)
add_textbox(slide, "Computer Science & Engineering  |  2024–25",
            0.8, 6.88, 11.7, 0.35,
            font_size=13, color=RGBColor(0x78, 0x9A, 0xBF), align=PP_ALIGN.CENTER)


# ── Save ────────────────────────────────────────────────────────────────────
out_path = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                        "Major_Project_Final_Presentation.pptx")
prs.save(out_path)
print(f"Saved: {out_path}")
print(f"Total slides: {len(prs.slides)}")
