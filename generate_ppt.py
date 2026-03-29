from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os
from PIL import Image

# ── Colour palette ──────────────────────────────────────────────────────────
DARK_BLUE   = RGBColor(0x0D, 0x2C, 0x54)   # deep navy  (backgrounds / headers)
MID_BLUE    = RGBColor(0x19, 0x76, 0xD2)   # primary blue
LIGHT_BLUE  = RGBColor(0xE3, 0xF2, 0xFD)   # pale blue  (content area bg)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT      = RGBColor(0x00, 0xB0, 0x8A)   # teal green (bullets / highlights)
DARK_TEXT   = RGBColor(0x1A, 0x1A, 0x2E)
GREY_TEXT   = RGBColor(0x55, 0x55, 0x66)
TICK_GREEN  = RGBColor(0x2E, 0x7D, 0x32)
PENDING_ORG = RGBColor(0xE6, 0x5C, 0x00)
CROSS_RED   = RGBColor(0xC6, 0x28, 0x28)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

BLANK = prs.slide_layouts[6]   # completely blank


# ── Helpers ──────────────────────────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill_rgb=None, line_rgb=None, line_width_pt=0):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.fill.background()
    if fill_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
    else:
        shape.fill.background()
    if line_rgb:
        shape.line.color.rgb = line_rgb
        shape.line.width = Pt(line_width_pt)
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, text, l, t, w, h,
                font_size=18, bold=False, color=WHITE,
                align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def header_bar(slide, title, subtitle=None):
    """Dark-blue top header bar."""
    add_rect(slide, 0, 0, 13.33, 1.3, fill_rgb=DARK_BLUE)
    add_textbox(slide, title, 0.35, 0.12, 12.0, 0.7,
                font_size=32, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_textbox(slide, subtitle, 0.35, 0.72, 12.0, 0.45,
                    font_size=16, color=RGBColor(0xBB, 0xDE, 0xFB), align=PP_ALIGN.LEFT)


def accent_line(slide):
    """Thin teal line below the header."""
    add_rect(slide, 0, 1.3, 13.33, 0.05, fill_rgb=ACCENT)


def slide_background(slide):
    """Light-blue full-slide background."""
    add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=LIGHT_BLUE)


def bullet_block(slide, items, l, t, w, col=DARK_TEXT,
                 size=17, marker="•", line_gap=0.38, bold_first=False):
    """Render a list of strings as bullet points."""
    for i, item in enumerate(items):
        txb = slide.shapes.add_textbox(
            Inches(l), Inches(t + i * line_gap), Inches(w), Inches(0.4))
        txb.word_wrap = True
        tf = txb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = f"{marker}  {item}"
        run.font.size = Pt(size)
        run.font.color.rgb = col
        run.font.bold = (bold_first and i == 0)


def add_image(slide, path, l, t, w, h):
    if os.path.exists(path):
        slide.shapes.add_picture(path, Inches(l), Inches(t), Inches(w), Inches(h))
    else:
        add_rect(slide, l, t, w, h, fill_rgb=RGBColor(0xCC, 0xCC, 0xCC))
        add_textbox(slide, "[Screenshot]", l + 0.1, t + h / 2 - 0.2, w - 0.2, 0.4,
                    font_size=14, color=DARK_TEXT, align=PP_ALIGN.CENTER)


def footer(slide, text="Smart Healthcare Appointment System  |  Mannat Vij"):
    add_rect(slide, 0, 7.18, 13.33, 0.32, fill_rgb=DARK_BLUE)
    add_textbox(slide, text, 0.3, 7.19, 12.7, 0.28,
                font_size=11, color=RGBColor(0xBB, 0xDE, 0xFB), align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 1 – Title
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)

# Full dark-blue background
add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=DARK_BLUE)

# Decorative teal block (left accent)
add_rect(slide, 0, 0, 0.12, 7.5, fill_rgb=ACCENT)

# Decorative teal block (bottom)
add_rect(slide, 0, 6.8, 13.33, 0.7, fill_rgb=RGBColor(0x0A, 0x1F, 0x3D))

# Subtle mid-tone rectangle behind title
add_rect(slide, 0.5, 1.6, 12.3, 3.5, fill_rgb=RGBColor(0x12, 0x3A, 0x6B))

# Title
add_textbox(slide, "Smart Healthcare", 0.8, 1.75, 11.7, 0.9,
            font_size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(slide, "Appointment System", 0.8, 2.55, 11.7, 0.9,
            font_size=52, bold=True, color=RGBColor(0x64, 0xB5, 0xF6), align=PP_ALIGN.CENTER)

# Teal separator line
add_rect(slide, 3.5, 3.55, 6.3, 0.06, fill_rgb=ACCENT)

add_textbox(slide, "MAJOR PROJECT PRESENTATION", 0.8, 3.72, 11.7, 0.5,
            font_size=18, bold=True, color=RGBColor(0xBB, 0xDE, 0xFB),
            align=PP_ALIGN.CENTER)

add_textbox(slide, "Presented by:  Mannat Vij", 0.8, 4.45, 11.7, 0.45,
            font_size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(slide, "Under the Guidance of:  Dr. Krishna Kumar", 0.8, 4.95, 11.7, 0.45,
            font_size=18, color=RGBColor(0xBB, 0xDE, 0xFB), align=PP_ALIGN.CENTER)

add_textbox(slide, "Computer Science & Engineering  |  2024–25",
            0.8, 6.88, 11.7, 0.35,
            font_size=13, color=RGBColor(0x78, 0x9A, 0xBF), align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 2 – Index
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
slide_background(slide)
header_bar(slide, "Index", "Presentation Overview")
accent_line(slide)
footer(slide)

index_items = [
    ("01", "Introduction",          "What the project is & why it matters"),
    ("02", "Objectives",            "Goals and purpose of the system"),
    ("03", "Methodology",           "Development approach, APIs & tools used"),
    ("04", "Progress & Status",     "Planned vs. completed features"),
    ("05", "Project Snapshots",     "Live screenshots from the running application"),
    ("06", "Future Scope",          "Planned enhancements and extensions"),
]

cols = [(0.45, 6.3), (6.9, 6.3)]
for idx, (num, title, desc) in enumerate(index_items):
    row = idx % 3
    col = idx // 3
    lx, lw = cols[col]
    ty = 1.55 + row * 1.7

    add_rect(slide, lx, ty, lw, 1.45,
             fill_rgb=WHITE,
             line_rgb=MID_BLUE, line_width_pt=1.2)

    # Number badge
    add_rect(slide, lx + 0.12, ty + 0.18, 0.55, 0.55, fill_rgb=MID_BLUE)
    add_textbox(slide, num, lx + 0.12, ty + 0.18, 0.55, 0.55,
                font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_textbox(slide, title, lx + 0.82, ty + 0.12, lw - 1.0, 0.42,
                font_size=20, bold=True, color=DARK_BLUE)
    add_textbox(slide, desc, lx + 0.82, ty + 0.55, lw - 1.0, 0.6,
                font_size=14, color=GREY_TEXT)


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 3 – Introduction
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
slide_background(slide)
header_bar(slide, "Introduction", "About the Project")
accent_line(slide)
footer(slide)

intro_para = (
    "In today's fast-paced world, the healthcare sector faces a critical challenge — "
    "the growing gap between patients seeking timely medical care and doctors managing "
    "overwhelming schedules through outdated, manual processes. Phone-based bookings, "
    "long waiting queues, and fragmented paper records result in missed consultations "
    "and poor patient experiences.\n\n"
    "The Smart Healthcare Appointment System is a robust, full-stack web application "
    "engineered to digitise and streamline the complete appointment lifecycle. Built on "
    "a React.js frontend (Material-UI), a Spring Boot RESTful backend, and MongoDB as a "
    "flexible NoSQL data store, the platform empowers patients to discover specialist "
    "doctors, check real-time availability, and book appointments instantly — while "
    "giving doctors full control over their schedules and patient queues.\n\n"
    "Secured by industry-standard JWT-based authentication and role-based access control, "
    "the system represents a scalable, modern solution to one of healthcare's most "
    "persistent operational challenges."
)

# White content card
add_rect(slide, 0.5, 1.55, 12.33, 5.4, fill_rgb=WHITE,
         line_rgb=RGBColor(0xC5, 0xD8, 0xF0), line_width_pt=1)

# Left accent stripe
add_rect(slide, 0.5, 1.55, 0.18, 5.4, fill_rgb=MID_BLUE)

txb = slide.shapes.add_textbox(Inches(0.9), Inches(1.75), Inches(11.8), Inches(5.0))
txb.word_wrap = True
tf = txb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.LEFT
run = p.add_run()
run.text = intro_para
run.font.size = Pt(17.5)
run.font.color.rgb = DARK_TEXT

# Tech-stack chips at bottom
chips = ["React.js", "Spring Boot 3", "MongoDB", "JWT Auth", "Material-UI", "REST API"]
chip_x = 0.9
for chip in chips:
    cw = len(chip) * 0.115 + 0.4
    add_rect(slide, chip_x, 6.55, cw, 0.38, fill_rgb=MID_BLUE)
    add_textbox(slide, chip, chip_x + 0.05, 6.56, cw - 0.1, 0.35,
                font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    chip_x += cw + 0.18


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 4 – Objectives
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
slide_background(slide)
header_bar(slide, "Objectives", "Goals & Purpose of the System")
accent_line(slide)
footer(slide)

objectives = [
    ("01", "Digitise Healthcare Access",
     "Design a secure, role-based platform (Patient / Doctor / Admin) that eliminates "
     "manual appointment booking and centralises healthcare management in a single portal."),
    ("02", "Scalable Backend Architecture",
     "Develop a production-grade RESTful API using Spring Boot 3 and MongoDB, with "
     "Swagger UI for centralised, interactive API documentation and endpoint management."),
    ("03", "Industry-Standard Security",
     "Implement JWT-based stateless authentication (HMAC-SHA256, 24-hr tokens) with "
     "BCrypt password hashing and Spring Security role guards to ensure data confidentiality."),
    ("04", "Intelligent Doctor Discovery",
     "Enable real-time doctor search with specialisation-based filtering, paginated results, "
     "live availability slots, and a visual time-slot booking interface for patients."),
    ("05", "End-to-End Appointment Lifecycle",
     "Allow doctors to manage consultation queues (Pending → Confirmed → Completed / "
     "Cancelled), update schedules dynamically, and maintain professional profiles."),
]

for i, (num, title, desc) in enumerate(objectives):
    ty = 1.55 + i * 1.08
    add_rect(slide, 0.45, ty, 12.4, 0.98,
             fill_rgb=WHITE, line_rgb=RGBColor(0xC5, 0xD8, 0xF0), line_width_pt=1)
    # Numbered badge
    add_rect(slide, 0.45, ty, 0.72, 0.98, fill_rgb=MID_BLUE)
    add_textbox(slide, num, 0.45, ty + 0.22, 0.72, 0.5,
                font_size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # Title
    add_textbox(slide, title, 1.3, ty + 0.05, 11.2, 0.38,
                font_size=17, bold=True, color=DARK_BLUE)
    # Description
    add_textbox(slide, desc, 1.3, ty + 0.44, 11.2, 0.5,
                font_size=13.5, color=GREY_TEXT)


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 5 – Methodology
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
slide_background(slide)
header_bar(slide, "Methodology", "Development Approach, APIs & Tools Used")
accent_line(slide)
footer(slide)

steps = [
    ("1", "Requirement Analysis & Planning",
     "Identified 3 user roles (Patient, Doctor, Admin), defined use cases, "
     "system architecture, API contracts, and data models."),
    ("2", "Environment Setup",
     "Java 17 · Spring Boot 3.2.3 · React 19 + Material-UI · MongoDB · Maven · Node.js 18"),
    ("3", "Database Design",
     "MongoDB collections: users (polymorphic), appointments, chat_sessions; "
     "Indexes: patientId, doctorId, email, username."),
    ("4", "Backend API Development  (Spring Boot REST)",
     "Auth: POST /api/auth/register, /api/auth/login  |  "
     "Users: GET/PUT /api/users/profile, PUT /api/users/password  |  "
     "Doctors: GET /api/doctors, /api/doctors/search, /api/doctors/available, "
     "PUT /api/doctors/availability  |  "
     "Appointments: POST /api/appointments, GET /api/appointments, "
     "PUT /api/appointments/{id}/status, DELETE /api/appointments/{id}"),
    ("5", "API Documentation — Swagger UI",
     "Integrated Springdoc OpenAPI; all endpoints centralised & testable at "
     "/swagger-ui.html — served as the single source of truth throughout development."),
    ("6", "Security Layer",
     "JWT (HMAC-SHA256, 24 hr expiry) · BCrypt password hashing · "
     "Spring Security filter chain · @PreAuthorize role-based guards per endpoint."),
    ("7", "Frontend Development",
     "React SPA with protected routes · Axios JWT interceptors · "
     "Role-aware dashboards · Responsive layout (desktop sidebar + mobile bottom nav)."),
    ("8", "Testing at Each Stage",
     "Postman collection for manual API testing · Swagger UI for live verification · "
     "Playwright E2E test suite (153 test cases) for full-flow validation."),
]

col1 = steps[:4]
col2 = steps[4:]

for col_idx, col_steps in enumerate([ col1, col2 ]):
    lx = 0.45 if col_idx == 0 else 6.88
    for row_i, (num, title, desc) in enumerate(col_steps):
        ty = 1.52 + row_i * 1.44
        add_rect(slide, lx, ty, 6.2, 1.34,
                 fill_rgb=WHITE, line_rgb=RGBColor(0xC5, 0xD8, 0xF0), line_width_pt=1)
        add_rect(slide, lx, ty, 0.48, 1.34, fill_rgb=ACCENT)
        add_textbox(slide, num, lx, ty + 0.42, 0.48, 0.45,
                    font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, title, lx + 0.58, ty + 0.05, 5.5, 0.4,
                    font_size=14, bold=True, color=DARK_BLUE)
        add_textbox(slide, desc, lx + 0.58, ty + 0.44, 5.5, 0.82,
                    font_size=11.5, color=GREY_TEXT)


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 6 – Progress
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
slide_background(slide)
header_bar(slide, "Progress & Current Status", "What Was Planned vs. What Is Done")
accent_line(slide)
footer(slide)

planned = [
    "User registration & login",
    "Doctor search & specialisation filter",
    "Appointment booking with time-slot picker",
    "Doctor availability management",
    "Role-based dashboards (Patient & Doctor)",
    "Profile management (edit + password change)",
    "Appointment status lifecycle",
    "Responsive UI (desktop + mobile)",
    "AI chatbot triage",
    "SMS / Email notifications",
    "Payment gateway integration",
]

done = [
    ("✔", "JWT-based auth (Patient & Doctor)",         TICK_GREEN),
    ("✔", "Search by name + filter by specialisation",  TICK_GREEN),
    ("✔", "Visual slot picker + booking form",          TICK_GREEN),
    ("✔", "Add / remove slots dynamically",             TICK_GREEN),
    ("✔", "Separate Patient & Doctor dashboards",       TICK_GREEN),
    ("✔", "Edit profile + change password",             TICK_GREEN),
    ("✔", "Pending → Confirmed → Completed / Cancelled", TICK_GREEN),
    ("✔", "Desktop sidebar + mobile bottom nav",        TICK_GREEN),
    ("~", "Data models ready, integration pending",     PENDING_ORG),
    ("✗", "Planned for next phase",                     CROSS_RED),
    ("✗", "Planned for next phase",                     CROSS_RED),
]

# Planned box
add_rect(slide, 0.45, 1.5, 5.9, 5.55,
         fill_rgb=WHITE, line_rgb=MID_BLUE, line_width_pt=1.5)
add_rect(slide, 0.45, 1.5, 5.9, 0.5, fill_rgb=DARK_BLUE)
add_textbox(slide, "PLANNED", 0.45, 1.5, 5.9, 0.5,
            font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

for i, item in enumerate(planned):
    add_textbox(slide, f"  ›  {item}", 0.55, 2.1 + i * 0.44, 5.7, 0.42,
                font_size=13, color=DARK_TEXT)

# Done box
add_rect(slide, 6.95, 1.5, 5.9, 5.55,
         fill_rgb=WHITE, line_rgb=ACCENT, line_width_pt=1.5)
add_rect(slide, 6.95, 1.5, 5.9, 0.5, fill_rgb=ACCENT)
add_textbox(slide, "COMPLETED / STATUS", 6.95, 1.5, 5.9, 0.5,
            font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

for i, (sym, text, col) in enumerate(done):
    add_textbox(slide, f"  {sym}  {text}", 7.05, 2.1 + i * 0.44, 5.7, 0.42,
                font_size=13, color=col)


# ════════════════════════════════════════════════════════════════════════════
#  SLIDES 7-10 – Project Snapshots
# ════════════════════════════════════════════════════════════════════════════

snapshot_slides = [
    ("Login & Registration", [
        ("screenshot_login.png",    0.45, 1.45, 5.9,  5.5,  "Login Page"),
        ("screenshot_register.png", 6.95, 1.45, 5.9,  5.5,  "Register Page"),
    ]),
    ("Patient Dashboard & Doctor Discovery", [
        ("screenshot_patient_dashboard.png", 0.45, 1.45, 5.9, 5.5, "Patient Dashboard"),
        ("screenshot_doctor_list.png",       6.95, 1.45, 5.9, 5.5, "Browse Doctors"),
    ]),
    ("Appointment Booking & Doctor Dashboard", [
        ("screenshot_book_appointment.png",  0.45, 1.45, 5.9, 5.5, "Book Appointment"),
        ("screenshot_doctor_dashboard.png",  6.95, 1.45, 5.9, 5.5, "Doctor Dashboard"),
    ]),
    ("Doctor Management & Swagger API Docs", [
        ("screenshot_doctor_appointments.png", 0.45, 1.45, 5.9, 5.5, "Doctor Appointments"),
        ("screenshot_swagger.png",             6.95, 1.45, 5.9, 5.5, "Swagger UI — API Docs"),
    ]),
]

for s_title, images in snapshot_slides:
    slide = prs.slides.add_slide(BLANK)
    slide_background(slide)
    header_bar(slide, "Project Snapshots", s_title)
    accent_line(slide)
    footer(slide)

    for (path, l, t, w, h, caption) in images:
        # Image frame border
        add_rect(slide, l - 0.05, t - 0.05, w + 0.1, h + 0.1,
                 fill_rgb=WHITE, line_rgb=MID_BLUE, line_width_pt=1.5)
        add_image(slide, path, l, t, w, h)
        # Caption bar
        add_rect(slide, l, t + h, w, 0.35, fill_rgb=DARK_BLUE)
        add_textbox(slide, caption, l, t + h, w, 0.35,
                    font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 11 – Future Scope
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
slide_background(slide)
header_bar(slide, "Future Scope", "Planned Enhancements & Extensions")
accent_line(slide)
footer(slide)

future_items = [
    ("SMS & Email Notifications",
     "Automated alerts to patients & doctors on booking, confirmation, and cancellation via SMS and email."),
    ("AI Symptom Checker Chatbot",
     "Patient describes symptoms; AI engine recommends the appropriate doctor specialisation before booking."),
    ("Prescription Management Portal",
     "Doctors upload digital prescriptions post-consultation; patients access complete prescription history in one place."),
    ("Online Payment Gateway",
     "Razorpay / Stripe integration for secure online consultation fee collection before appointments."),
    ("Video Consultation (WebRTC)",
     "Built-in virtual appointment feature for remote consultations without leaving the platform."),
    ("Calendar Sync & Reminders",
     "Google Calendar / iCal export so patients and doctors receive appointment reminders on their devices."),
    ("Admin Analytics Dashboard",
     "Visual reports on appointment trends, doctor performance metrics, and user growth statistics."),
]

cols_fs = [future_items[:4], future_items[4:]]
col_xs  = [0.45, 7.0]
col_w   = 6.0

for ci, col_items in enumerate(cols_fs):
    lx = col_xs[ci]
    for ri, (title, desc) in enumerate(col_items):
        ty = 1.55 + ri * 1.44
        add_rect(slide, lx, ty, col_w, 1.3,
                 fill_rgb=WHITE, line_rgb=RGBColor(0xC5, 0xD8, 0xF0), line_width_pt=1)
        add_rect(slide, lx, ty, 0.15, 1.3, fill_rgb=ACCENT)
        add_textbox(slide, title, lx + 0.28, ty + 0.08, col_w - 0.35, 0.38,
                    font_size=15, bold=True, color=DARK_BLUE)
        add_textbox(slide, desc, lx + 0.28, ty + 0.48, col_w - 0.35, 0.72,
                    font_size=12.5, color=GREY_TEXT)


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 12 – Thank You
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=DARK_BLUE)
add_rect(slide, 0, 0, 0.12, 7.5, fill_rgb=ACCENT)
add_rect(slide, 0, 6.8, 13.33, 0.7, fill_rgb=RGBColor(0x0A, 0x1F, 0x3D))
add_rect(slide, 0.5, 2.2, 12.3, 2.8, fill_rgb=RGBColor(0x12, 0x3A, 0x6B))

add_textbox(slide, "Thank You", 0.8, 2.4, 11.7, 1.1,
            font_size=60, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_rect(slide, 4.0, 3.55, 5.3, 0.07, fill_rgb=ACCENT)
add_textbox(slide, "Smart Healthcare Appointment System", 0.8, 3.75, 11.7, 0.5,
            font_size=20, color=RGBColor(0xBB, 0xDE, 0xFB), align=PP_ALIGN.CENTER)
add_textbox(slide, "Mannat Vij  |  Under the Guidance of Dr. Krishna Kumar",
            0.8, 4.3, 11.7, 0.4,
            font_size=16, color=RGBColor(0x9E, 0xC8, 0xF0), align=PP_ALIGN.CENTER)
add_textbox(slide, "Computer Science & Engineering  |  2024–25",
            0.8, 6.88, 11.7, 0.35,
            font_size=13, color=RGBColor(0x78, 0x9A, 0xBF), align=PP_ALIGN.CENTER)


# ── Save ─────────────────────────────────────────────────────────────────────
out_path = r"C:\Users\manna\IdeaProjects\Major_Project\Major_Project_Presentation.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
print(f"Total slides: {len(prs.slides)}")
